"""文档加载器 — 解析多种文件格式为纯文本，供 RAG 管道入库

支持格式：
- .txt       : 纯文本
- .md        : Markdown
- .docx      : Word 文档
- .pdf       : PDF 文档
- .pptx      : PowerPoint 演示文稿
"""

from pathlib import Path


def load_document(file_path: str | Path) -> str:
    """根据文件后缀自动选择合适的解析器，返回纯文本内容"""
    file_path = Path(file_path)
    suffix = file_path.suffix.lower()

    if suffix == ".txt":
        return _load_txt(file_path)
    elif suffix == ".md":
        return _load_txt(file_path)  # Markdown 本质是纯文本
    elif suffix == ".docx":
        return _load_docx(file_path)
    elif suffix == ".pdf":
        return _load_pdf(file_path)
    elif suffix == ".pptx":
        return _load_pptx(file_path)
    else:
        raise ValueError(f"不支持的文件格式: {suffix}")


# ========== 各格式解析器 ==========

def _load_txt(file_path: Path) -> str:
    """纯文本 / Markdown"""
    return file_path.read_text(encoding="utf-8")


def _load_docx(file_path: Path) -> str:
    """Word 文档 — 逐段提取文字"""
    from docx import Document

    doc = Document(str(file_path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _load_pdf(file_path: Path) -> str:
    """PDF 文档 — 逐页提取文字"""
    import logging
    import pdfplumber

    # 抑制 pdfplumber 的 FontBBox 无害警告
    logging.getLogger("pdfplumber").setLevel(logging.ERROR)

    with pdfplumber.open(str(file_path)) as pdf:
        pages = [page.extract_text() or "" for page in pdf.pages]
    return "\n\n".join(pages)


def _load_pptx(file_path: Path) -> str:
    """PPT 演示文稿 — 逐页逐形状提取文字"""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n\n".join(texts)
